from __future__ import annotations

from dataclasses import dataclass, field
import math
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("docs/game_ui_windows_interaction_points.pptx")

SLIDE_W = 13.333333
SLIDE_H = 7.5
GAME_VIEWPORT_W = 1920.0
GAME_VIEWPORT_H = 1080.0
WINDOW_OVERRIDE_W = 960
WINDOW_OVERRIDE_H = 540
GAME_FRAME_H = 5.95
GAME_FRAME_W = GAME_FRAME_H * GAME_VIEWPORT_W / GAME_VIEWPORT_H
GAME_FRAME_X = (SLIDE_W - GAME_FRAME_W) / 2
GAME_FRAME_Y = 0.78

ROOT = Path(__file__).resolve().parents[1]

COLORS = {
    "bg": RGBColor(14, 24, 39),
    "panel": RGBColor(30, 45, 68),
    "panel2": RGBColor(42, 59, 84),
    "button": RGBColor(74, 144, 226),
    "button2": RGBColor(95, 177, 234),
    "input": RGBColor(104, 94, 180),
    "drag": RGBColor(46, 170, 132),
    "display": RGBColor(90, 103, 122),
    "text": RGBColor(238, 244, 255),
    "muted": RGBColor(176, 190, 210),
    "warn": RGBColor(245, 180, 65),
    "frame": RGBColor(255, 225, 130),
}

CONTROL_TYPES = {
    "Button",
    "CheckBox",
    "ColorRect",
    "Control",
    "HBoxContainer",
    "Label",
    "LineEdit",
    "MarginContainer",
    "NinePatchRect",
    "Panel",
    "ProgressBar",
    "RichTextLabel",
    "ScrollContainer",
    "TextureButton",
    "TextureProgressBar",
    "TextureRect",
    "TextEdit",
    "VBoxContainer",
}

BUTTON_TYPES = {"Button", "TextureButton"}
INPUT_TYPES = {"LineEdit", "TextEdit"}


@dataclass
class SceneNode:
    name: str
    type_name: str
    path: str
    parent_path: str | None
    props: dict[str, object] = field(default_factory=dict)


@dataclass
class LayoutNode:
    node: SceneNode
    rect: tuple[float, float, float, float] | None
    origin: tuple[float, float]
    scale: tuple[float, float]
    size: tuple[float, float]


def inch(value: float):
    return Inches(value)


def px_x(value: float) -> float:
    return GAME_FRAME_X + value / GAME_VIEWPORT_W * GAME_FRAME_W


def px_y(value: float) -> float:
    return GAME_FRAME_Y + value / GAME_VIEWPORT_H * GAME_FRAME_H


def px_w(value: float) -> float:
    return value / GAME_VIEWPORT_W * GAME_FRAME_W


def px_h(value: float) -> float:
    return value / GAME_VIEWPORT_H * GAME_FRAME_H


def clamp_text_size(height_px: float, base: int = 9) -> int:
    if height_px < 24:
        return 5
    if height_px < 36:
        return 6
    if height_px < 50:
        return 7
    if height_px < 80:
        return min(base, 9)
    return base


def add_textbox(slide, x, y, w, h, text, size=14, color="text", bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = inch(0.06)
    tf.margin_right = inch(0.06)
    tf.margin_top = inch(0.02)
    tf.margin_bottom = inch(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COLORS[color]
    return box


def add_rect(slide, x, y, w, h, text="", fill="panel", line="muted", size=12, bold=False):
    if abs(w) < 0.01 or abs(h) < 0.01:
        return None
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS[fill]
    shape.line.color.rgb = COLORS[line]
    shape.line.width = Pt(1.1)
    if text:
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = inch(0.03)
        tf.margin_right = inch(0.03)
        tf.margin_top = inch(0.02)
        tf.margin_bottom = inch(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = COLORS["text"]
    return shape


def add_button(slide, x, y, w, h, text, size=12):
    return add_rect(slide, x, y, w, h, text, fill="button", line="button2", size=size, bold=True)


def add_input(slide, x, y, w, h, text, size=12):
    return add_rect(slide, x, y, w, h, text, fill="input", line="button2", size=size, bold=True)


def add_drag(slide, x, y, w, h, text, size=12):
    return add_rect(slide, x, y, w, h, text, fill="drag", line="button2", size=size, bold=True)


def add_display(slide, x, y, w, h, text, size=11):
    return add_rect(slide, x, y, w, h, text, fill="display", line="muted", size=size)


def add_rect_px(slide, left, top, right, bottom, text="", fill="panel", line="muted", size=9, bold=False):
    return add_rect(
        slide,
        px_x(left),
        px_y(top),
        px_w(right - left),
        px_h(bottom - top),
        text,
        fill=fill,
        line=line,
        size=size,
        bold=bold,
    )


def add_button_px(slide, left, top, right, bottom, text, size=9):
    return add_rect_px(slide, left, top, right, bottom, text, fill="button", line="button2", size=size, bold=True)


def add_input_px(slide, left, top, right, bottom, text, size=9):
    return add_rect_px(slide, left, top, right, bottom, text, fill="input", line="button2", size=size, bold=True)


def add_drag_px(slide, left, top, right, bottom, text, size=9):
    return add_rect_px(slide, left, top, right, bottom, text, fill="drag", line="button2", size=size, bold=True)


def add_display_px(slide, left, top, right, bottom, text, size=8):
    return add_rect_px(slide, left, top, right, bottom, text, fill="display", line="muted", size=size)


def add_header(slide, title, scene_path, count_text):
    add_textbox(slide, 0.45, 0.18, 8.8, 0.35, title, size=20, bold=True)
    add_textbox(slide, 0.47, 0.56, 8.7, 0.25, scene_path, size=8, color="muted")
    add_rect(slide, 10.4, 0.22, 2.45, 0.46, count_text, fill="panel2", line="warn", size=11, bold=True)


def add_game_window_guide(slide):
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        inch(GAME_FRAME_X),
        inch(GAME_FRAME_Y),
        inch(GAME_FRAME_W),
        inch(GAME_FRAME_H),
    )
    frame.fill.background()
    frame.line.color.rgb = COLORS["frame"]
    frame.line.width = Pt(2)
    add_textbox(
        slide,
        GAME_FRAME_X,
        GAME_FRAME_Y + GAME_FRAME_H + 0.02,
        5.0,
        0.18,
        f"游戏视口：{int(GAME_VIEWPORT_W)}x{int(GAME_VIEWPORT_H)}；窗口覆盖：{WINDOW_OVERRIDE_W}x{WINDOW_OVERRIDE_H}",
        size=7,
        color="frame",
    )


def add_legend(slide):
    add_button(slide, 9.75, 6.85, 0.55, 0.23, "", size=8)
    add_textbox(slide, 10.35, 6.83, 0.8, 0.22, "按钮", size=8, color="muted")
    add_input(slide, 11.1, 6.85, 0.55, 0.23, "", size=8)
    add_textbox(slide, 11.7, 6.83, 0.7, 0.22, "输入", size=8, color="muted")
    add_drag(slide, 12.25, 6.85, 0.55, 0.23, "", size=8)
    add_textbox(slide, 12.85, 6.83, 0.45, 0.22, "拖拽", size=8, color="muted")


def new_slide(prs, title, scene_path, count_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["bg"]
    add_header(slide, title, scene_path, count_text)
    add_legend(slide)
    return slide


def parse_attrs(raw: str) -> dict[str, str]:
    attrs: dict[str, str] = {}
    for key, quoted, bare in re.findall(r"(\w+)=(?:\"([^\"]*)\"|([^\s\]]+))", raw):
        attrs[key] = quoted if quoted != "" else bare
    return attrs


def parse_value(raw: str):
    raw = raw.strip()
    if raw.startswith('"') and raw.endswith('"'):
        return raw[1:-1]
    if raw == "true":
        return True
    if raw == "false":
        return False
    vector = re.match(r"Vector2\(([-0-9.eE]+),\s*([-0-9.eE]+)\)", raw)
    if vector:
        return (float(vector.group(1)), float(vector.group(2)))
    try:
        return float(raw)
    except ValueError:
        return raw


def scene_file(scene_path: str) -> Path:
    return ROOT / scene_path


def parse_tscn(scene_path: str) -> dict[str, SceneNode]:
    nodes: dict[str, SceneNode] = {}
    current: SceneNode | None = None
    root_name: str | None = None
    text = scene_file(scene_path).read_text(encoding="utf-8-sig", errors="replace")
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if line.startswith("[node "):
            attrs = parse_attrs(line)
            name = attrs["name"]
            type_name = attrs.get("type", "")
            parent = attrs.get("parent")
            if parent is None:
                root_name = name
                path = name
                parent_path = None
            else:
                assert root_name is not None
                parent_suffix = "" if parent == "." else "/" + parent
                path = f"{root_name}{parent_suffix}/{name}"
                parent_path = root_name if parent == "." else f"{root_name}/{parent}"
            current = SceneNode(name, type_name, path, parent_path)
            nodes[path] = current
            continue
        if current is None or "=" not in line or line.startswith("["):
            continue
        key, value = line.split("=", 1)
        current.props[key.strip()] = parse_value(value)
    return nodes


def prop_float(node: SceneNode, key: str, default: float = 0.0) -> float:
    value = node.props.get(key, default)
    return float(value) if isinstance(value, (int, float)) else default


def prop_vec2(node: SceneNode, key: str, default: tuple[float, float]) -> tuple[float, float]:
    value = node.props.get(key, default)
    return value if isinstance(value, tuple) else default


def compute_layout(scene_path: str) -> dict[str, LayoutNode]:
    nodes = parse_tscn(scene_path)
    layout: dict[str, LayoutNode] = {}

    for path, node in nodes.items():
        if node.parent_path and node.parent_path in layout:
            parent = layout[node.parent_path]
        else:
            parent = LayoutNode(
                node=node,
                rect=(0.0, 0.0, GAME_VIEWPORT_W, GAME_VIEWPORT_H),
                origin=(0.0, 0.0),
                scale=(1.0, 1.0),
                size=(GAME_VIEWPORT_W, GAME_VIEWPORT_H),
            )

        parent_origin = parent.origin
        parent_scale = parent.scale
        parent_size = parent.size
        node_scale = prop_vec2(node, "scale", (1.0, 1.0))

        if node.type_name in CONTROL_TYPES:
            anchor_left = prop_float(node, "anchor_left")
            anchor_top = prop_float(node, "anchor_top")
            anchor_right = prop_float(node, "anchor_right")
            anchor_bottom = prop_float(node, "anchor_bottom")
            left = anchor_left * parent_size[0] + prop_float(node, "offset_left")
            top = anchor_top * parent_size[1] + prop_float(node, "offset_top")
            right = anchor_right * parent_size[0] + prop_float(node, "offset_right")
            bottom = anchor_bottom * parent_size[1] + prop_float(node, "offset_bottom")
            width = right - left
            height = bottom - top
            min_size = prop_vec2(node, "custom_minimum_size", (0.0, 0.0))
            if width <= 0.0 and min_size[0] > 0.0:
                width = min_size[0]
                right = left + width
            if height <= 0.0 and min_size[1] > 0.0:
                height = min_size[1]
                bottom = top + height

            global_left = parent_origin[0] + left * parent_scale[0]
            global_top = parent_origin[1] + top * parent_scale[1]
            global_right = global_left + width * parent_scale[0] * node_scale[0]
            global_bottom = global_top + height * parent_scale[1] * node_scale[1]
            origin = (global_left, global_top)
            scale = (parent_scale[0] * node_scale[0], parent_scale[1] * node_scale[1])
            size = (width, height)
            rect = (global_left, global_top, global_right, global_bottom)
        elif node.type_name == "Node2D":
            position = prop_vec2(node, "position", (0.0, 0.0))
            origin = (
                parent_origin[0] + position[0] * parent_scale[0],
                parent_origin[1] + position[1] * parent_scale[1],
            )
            scale = (parent_scale[0] * node_scale[0], parent_scale[1] * node_scale[1])
            size = (0.0, 0.0)
            rect = None
        else:
            origin = parent_origin
            scale = (parent_scale[0] * node_scale[0], parent_scale[1] * node_scale[1])
            size = parent_size
            rect = None

        layout[path] = LayoutNode(node=node, rect=rect, origin=origin, scale=scale, size=size)

    return layout


def find_layout(layout: dict[str, LayoutNode], name: str) -> LayoutNode:
    matches = [item for path, item in layout.items() if path.endswith("/" + name) or item.node.name == name]
    if not matches:
        raise KeyError(name)
    return matches[0]


def rect_for(layout: dict[str, LayoutNode], name: str) -> tuple[float, float, float, float]:
    rect = find_layout(layout, name).rect
    if rect is None:
        raise KeyError(name)
    return rect


def draw_rect_from_layout(slide, layout: dict[str, LayoutNode], name: str, text: str, kind="display", size=8):
    left, top, right, bottom = rect_for(layout, name)
    if kind == "button":
        return add_button_px(slide, left, top, right, bottom, text, size=clamp_text_size(bottom - top, size))
    if kind == "input":
        return add_input_px(slide, left, top, right, bottom, text, size=clamp_text_size(bottom - top, size))
    if kind == "drag":
        return add_drag_px(slide, left, top, right, bottom, text, size=clamp_text_size(bottom - top, size))
    return add_display_px(slide, left, top, right, bottom, text, size=clamp_text_size(bottom - top, size))


def draw_buttons_from_layout(slide, layout: dict[str, LayoutNode], labels: dict[str, str] | None = None):
    labels = labels or {}
    for item in layout.values():
        if item.node.type_name not in BUTTON_TYPES or item.rect is None:
            continue
        label = labels.get(item.node.name, item.node.name)
        left, top, right, bottom = item.rect
        add_button_px(slide, left, top, right, bottom, label, size=clamp_text_size(bottom - top, 8))


def world_map_node_positions() -> list[tuple[int, float, float, float]]:
    result = [(0, 700.0, 590.0, 38.0)]
    node_id = 1
    for ring_index, (count, radius) in enumerate([(5, 200.0), (8, 330.0), (11, 455.0)]):
        angle_offset = -math.pi * 0.5 + float(ring_index) * 0.11
        for i in range(count):
            angle = angle_offset + math.tau * (float(i) + 0.5) / float(count)
            result.append((node_id, 700.0 + math.cos(angle) * radius, 590.0 + math.sin(angle) * radius, 24.0))
            node_id += 1
    return result


def add_mapping_note(slide, text: str):
    add_textbox(slide, GAME_FRAME_X + 5.05, GAME_FRAME_Y + GAME_FRAME_H + 0.02, 5.45, 0.18, text, size=7, color="muted")


def add_function_annotations(
    slide,
    title: str,
    scene_path: str,
    items: list[tuple[str, str]],
    extra: str = "",
):
    visible_lines = ["按钮功能"]
    visible_lines.extend(f"{name}：{desc}" for name, desc in items[:7])
    if len(items) > 7:
        visible_lines.append("更多见备注区")
    if not items:
        visible_lines.append("无按钮")
    add_rect(slide, 12.02, 1.05, 1.22, 5.52, "", fill="panel2", line="muted")
    add_textbox(slide, 12.06, 1.11, 1.14, 5.38, "\n".join(visible_lines), size=6, color="text")

    notes_lines = [title, f"场景：{scene_path}", "", "按钮/交互点功能："]
    if items:
        notes_lines.extend(f"- {name}：{desc}" for name, desc in items)
    else:
        notes_lines.append("- 无按钮：该页面只显示状态或装饰信息。")
    if extra:
        notes_lines.extend(["", "补充说明：", extra])
    slide.notes_slide.notes_text_frame.text = "\n".join(notes_lines)


def slide_main_menu(prs):
    scene = "scenes/ui/main_menu/MainMenuGeneratedUI.tscn"
    title = "主菜单"
    slide = new_slide(prs, title, scene, "5 个交互点")
    layout = compute_layout(scene)
    draw_rect_from_layout(slide, layout, "Background", "背景")
    labels = {
        "StartFrame": "开始游戏",
        "ExploreFrame": "探索房间",
        "BossFrame": "Boss 挑战",
        "SettingsFrame": "设置",
        "QuitFrame": "退出游戏",
    }
    for name, label in labels.items():
        draw_rect_from_layout(slide, layout, name, label, kind="button", size=10)
    add_mapping_note(slide, "位置来源：MainMenuGeneratedUI.tscn 的 Frame 矩形。")
    add_function_annotations(slide, title, scene, [
        ("开始游戏", "重置流程状态，开启正式流程并进入世界地图。"),
        ("探索房间", "取消正式流程，直接进入探索房间。"),
        ("Boss 挑战", "重置玩家状态并进入 Boss 选择界面。"),
        ("设置", "打开主菜单设置弹窗。"),
        ("退出游戏", "退出游戏程序。"),
    ], "按钮悬停动画由 scripts/app/MainMenu.gd 控制，PPT 中只表示默认状态下的实际位置。")


def slide_settings(prs):
    scene = "scenes/ui/main_menu/SettingsPopup.tscn"
    title = "设置弹窗"
    slide = new_slide(prs, title, scene, "1 个交互点")
    layout = compute_layout(scene)
    draw_rect_from_layout(slide, layout, "Shade", "遮罩")
    draw_rect_from_layout(slide, layout, "Panel", "面板")
    draw_rect_from_layout(slide, layout, "TitleLabel", "标题")
    draw_rect_from_layout(slide, layout, "SettingsMessage", "提示文本")
    draw_rect_from_layout(slide, layout, "CloseButton", "关闭", kind="button", size=10)
    add_mapping_note(slide, "位置来源：SettingsPopup.tscn 的全局 Control 矩形。")
    add_function_annotations(slide, title, scene, [
        ("关闭", "关闭设置弹窗并回到主菜单。"),
    ])


def slide_boss_select(prs):
    scene = "scenes/ui/boss_select/BossSelectUI.tscn"
    title = "Boss 选择"
    slide = new_slide(prs, title, scene, "21 个交互点")
    layout = compute_layout(scene)
    draw_rect_from_layout(slide, layout, "TitleLabel", "标题")
    labels = {f"Boss{i}Button": f"Boss {i}" for i in range(1, 21)}
    labels["BackButton"] = "返回"
    draw_buttons_from_layout(slide, layout, labels)
    add_mapping_note(slide, "位置来源：BossSelectUI.tscn 的 Button 矩形。")
    add_function_annotations(slide, title, scene, [
        ("Boss 1-20", "进入对应 Boss 战斗场景。具体场景映射在 scripts/app/BossSelect.gd。"),
        ("返回", "返回主菜单。"),
    ], "PPT 中用编号表示 Boss 按钮，便于编辑；实际按钮文本来自 BossSelectUI.tscn。")


def slide_game_over(prs):
    scene = "scenes/ui/game_over/GameOverUI.tscn"
    title = "游戏结束"
    slide = new_slide(prs, title, scene, "1 个交互点")
    layout = compute_layout(scene)
    draw_rect_from_layout(slide, layout, "Overlay", "遮罩")
    draw_rect_from_layout(slide, layout, "TitleLabel", "标题")
    draw_rect_from_layout(slide, layout, "FinalScoreLabel", "结算信息")
    draw_rect_from_layout(slide, layout, "RestartButton", "重新开始", kind="button", size=9)
    add_mapping_note(slide, "位置来源：GameOverUI.tscn 的全局 Control 矩形。")
    add_function_annotations(slide, title, scene, [
        ("重新开始", "重置游戏状态并返回主菜单或重新开始流程，逻辑在 scripts/app/GameOver.gd。"),
    ])


def slide_world_map(prs):
    scene = "scenes/ui/world_map/WorldMapUI.tscn"
    title = "世界地图"
    slide = new_slide(prs, title, scene, "4 个固定按钮 + 25 个节点点击")
    layout = compute_layout(scene)
    add_drag_px(slide, 0, 0, 1920, 1080, "地图绘制/节点点击区", size=11)
    draw_rect_from_layout(slide, layout, "TopBar", "顶部状态")
    draw_rect_from_layout(slide, layout, "TitleLabel", "标题")
    draw_rect_from_layout(slide, layout, "StatsLabel", "统计")
    draw_rect_from_layout(slide, layout, "DetailsPanel", "详情面板")
    draw_rect_from_layout(slide, layout, "DetailsBody", "详情文本")
    draw_rect_from_layout(slide, layout, "MessageLabel", "消息")
    draw_buttons_from_layout(slide, layout, {
        "EnterButton": "进入",
        "ShopButton": "商店",
        "HangarButton": "机库",
        "BackButton": "返回",
    })
    for node_id, x, y, radius in world_map_node_positions():
        label = "基地" if node_id == 0 else str(node_id)
        add_drag_px(slide, x - radius, y - radius, x + radius, y + radius, label, size=5 if node_id else 7)
    add_mapping_note(slide, "固定控件来自 WorldMapUI.tscn；节点来自 RunManager._generate_world_map()。")
    add_function_annotations(slide, title, scene, [
        ("进入", "进入当前选中的节点；基地警报时进入危机 Boss。"),
        ("商店", "在基地节点打开商店弹窗。"),
        ("机库", "在基地节点打开机库弹窗。"),
        ("返回", "返回主菜单。"),
        ("地图节点 1-24", "选中节点并刷新右侧详情；可进入节点由 RunManager 状态决定。"),
        ("基地", "选中方舟核心，可进入商店、机库或危机 Boss。"),
    ])


def draw_dynamic_item_rows(slide, scroll_rect: tuple[float, float, float, float], count: int, action_label: str):
    left, top, right, _bottom = scroll_rect
    row_h = 116.0
    gap = 12.0
    button_w = 130.0
    button_h = 48.0
    for i in range(count):
        row_top = top + i * (row_h + gap)
        row_bottom = row_top + row_h
        add_display_px(slide, left, row_top, right, row_bottom, f"动态物品行 {i + 1}", size=6)
        button_top = row_top + (row_h - button_h) * 0.5
        add_button_px(slide, right - button_w, button_top, right, button_top + button_h, action_label, size=6)


def slide_shop(prs):
    scene = "scenes/ui/world_map/ShopPopup.tscn"
    title = "商店弹窗"
    slide = new_slide(prs, title, scene, "1 个固定按钮 + 最多 8 个物品按钮")
    layout = compute_layout(scene)
    for name, label in [
        ("Shade", "遮罩"),
        ("Panel", "面板"),
        ("TitleLabel", "标题"),
        ("MineralsLabel", "星髓矿"),
        ("HintLabel", "提示"),
        ("ItemsScroll", "物品滚动区"),
        ("MessageLabel", "消息"),
    ]:
        draw_rect_from_layout(slide, layout, name, label)
    draw_rect_from_layout(slide, layout, "CloseButton", "关闭", kind="button", size=9)
    draw_dynamic_item_rows(slide, rect_for(layout, "ItemsScroll"), 8, "购买")
    add_mapping_note(slide, "静态控件来自 ShopPopup.tscn；物品行按 EquipmentItemRow 行高放入 ItemsScroll。")
    add_function_annotations(slide, title, scene, [
        ("购买", "购买对应装备；已拥有时按钮禁用。购买结果会刷新库存与提示消息。"),
        ("关闭", "关闭商店弹窗并回到世界地图。"),
    ], "商店物品行由 scripts/ui/world_map/ShopPopup.gd 在运行时根据 EquipmentCatalog 动态生成。")


def slide_hangar(prs):
    scene = "scenes/ui/world_map/HangarPopup.tscn"
    title = "机库弹窗"
    slide = new_slide(prs, title, scene, "1 个固定按钮 + 库存物品按钮")
    layout = compute_layout(scene)
    for name, label in [
        ("Shade", "遮罩"),
        ("Panel", "面板"),
        ("TitleLabel", "标题"),
        ("WeaponLabel", "当前武器"),
        ("AuxLabel", "辅助机"),
        ("ComputeLabel", "算力"),
        ("ItemsScroll", "物品滚动区"),
        ("MessageLabel", "消息"),
    ]:
        draw_rect_from_layout(slide, layout, name, label)
    draw_rect_from_layout(slide, layout, "CloseButton", "关闭", kind="button", size=9)
    draw_dynamic_item_rows(slide, rect_for(layout, "ItemsScroll"), 7, "装配")
    add_mapping_note(slide, "静态控件来自 HangarPopup.tscn；物品行按 EquipmentItemRow 行高放入 ItemsScroll。")
    add_function_annotations(slide, title, scene, [
        ("装配", "装配对应武器或辅助机；已装配武器按钮禁用。"),
        ("卸下", "已装配辅助机时，按钮会切换为卸下。"),
        ("关闭", "关闭机库弹窗并回到世界地图。"),
    ], "机库物品行由 scripts/ui/world_map/HangarPopup.gd 根据玩家库存动态生成。")


def slide_equipment_row(prs):
    scene = "scenes/ui/world_map/EquipmentItemRow.tscn"
    title = "装备物品行组件"
    slide = new_slide(prs, title, scene, "1 个交互点")
    # The row is instantiated into Shop/Hangar scroll containers. This page uses the
    # first shop row's runtime parent rect so its position is still viewport-mapped.
    shop_layout = compute_layout("scenes/ui/world_map/ShopPopup.tscn")
    left, top, right, _bottom = rect_for(shop_layout, "ItemsScroll")
    row_h = 116.0
    add_display_px(slide, left, top, right, top + row_h, "装备物品行运行时位置")
    add_button_px(slide, right - 130.0, top + 34.0, right, top + 82.0, "操作", size=7)
    add_mapping_note(slide, "组件按商店第一个物品行的运行时位置展示。")
    add_function_annotations(slide, title, scene, [
        ("操作", "由父弹窗决定功能：商店中为购买，机库中为装配或卸下。"),
    ], "这是复用组件，不是独立窗口；PPT 中放在商店 ItemsScroll 的第一行位置便于编辑。")


def slide_event_result(prs):
    scene = "scenes/ui/world_map/EventResultPopup.tscn"
    title = "事件结算弹窗"
    slide = new_slide(prs, title, scene, "1 个交互点")
    layout = compute_layout(scene)
    for name, label in [
        ("Shade", "遮罩"),
        ("Panel", "面板"),
        ("TitleLabel", "标题"),
        ("BodyLabel", "结算内容"),
    ]:
        draw_rect_from_layout(slide, layout, name, label)
    draw_rect_from_layout(slide, layout, "CloseButton", "确认", kind="button", size=9)
    add_mapping_note(slide, "位置来源：EventResultPopup.tscn 的全局 Control 矩形。")
    add_function_annotations(slide, title, scene, [
        ("确认", "关闭事件结算弹窗并返回世界地图。"),
    ])


def slide_evacuation(prs):
    scene = "scenes/ui/EvacuationSuccessHUD.tscn"
    title = "撤离成功"
    slide = new_slide(prs, title, scene, "1 个交互点")
    layout = compute_layout(scene)
    draw_rect_from_layout(slide, layout, "Panel", "面板")
    draw_rect_from_layout(slide, layout, "Title", "标题")
    draw_rect_from_layout(slide, layout, "EvacuateButton", "撤离", kind="button", size=9)
    add_mapping_note(slide, "位置来源：EvacuationSuccessHUD.tscn 的全局 Control 矩形。")
    add_function_annotations(slide, title, scene, [
        ("撤离", "确认撤离成功，触发探索房间的撤离回调。"),
    ])


def slide_command_console(prs):
    scene = "scenes/ui/explore/CommandConsolePopup.tscn"
    title = "命令台"
    slide = new_slide(prs, title, scene, "1 个交互点")
    layout = compute_layout(scene)
    draw_rect_from_layout(slide, layout, "CommandDialogPanel", "对话面板")
    draw_rect_from_layout(slide, layout, "CommandDialogLabel", "对话文本")
    draw_rect_from_layout(slide, layout, "CommandInputPanel", "输入面板")
    draw_rect_from_layout(slide, layout, "CommandInputEdit", "命令输入框", kind="input", size=9)
    add_mapping_note(slide, "位置来源：CommandConsolePopup.tscn；节点默认隐藏，打开命令台时显示。")
    add_function_annotations(slide, title, scene, [
        ("命令输入框", "输入调试/控制命令并提交。"),
    ], "这是输入交互，不是按钮；具体命令解析在探索房间相关脚本中。")


def slide_explore_map(prs):
    scene = "scenes/ui/explore/ExploreMapUI.tscn"
    title = "探索地图"
    slide = new_slide(prs, title, scene, "1 个拖拽交互")
    layout = compute_layout(scene)
    left, top, right, bottom = rect_for(layout, "ExploreMapUI")
    add_drag_px(slide, left, top, right, bottom, "探索地图拖拽区", size=12)
    add_mapping_note(slide, "位置来源：ExploreMapUI.tscn 根 Control 矩形。")
    add_function_annotations(slide, title, scene, [
        ("探索地图拖拽区", "按住并拖动地图，查看探索房间整体布局。"),
    ], "房间图标由运行时动态生成，图标本身忽略鼠标输入。")


def slide_compass(prs):
    scene = "scenes/ui/explore/CompassMiniMap.tscn"
    title = "罗盘小地图"
    slide = new_slide(prs, title, scene, "0 个交互点")
    layout = compute_layout(scene)
    draw_rect_from_layout(slide, layout, "CompassMiniMap", "罗盘小地图")
    add_mapping_note(slide, "位置来源：CompassMiniMap.tscn 根 Control 矩形。")
    add_function_annotations(slide, title, scene, [], "显示玩家周边、奖励、障碍等方向信息；不接收鼠标交互。")


def slide_loading(prs):
    scene = "scenes/ui/explore_loading/ExploreLoadingScreen.tscn"
    title = "探索加载界面"
    slide = new_slide(prs, title, scene, "0 个交互点")
    layout = compute_layout(scene)
    for name, label in [
        ("BlackOverlay", "黑色底层"),
        ("BlueWash", "蓝色叠层"),
        ("LogoPlate", "标题底板"),
        ("Title", "标题"),
        ("SpinnerPivot", "旋转标志"),
        ("LoadingBarArt", "进度条素材"),
        ("Panel", "进度面板"),
        ("ProgressBar", "进度条"),
        ("Label", "百分比"),
        ("TipPanel", "提示底板"),
        ("TipLabel", "提示文本"),
    ]:
        draw_rect_from_layout(slide, layout, name, label)
    add_mapping_note(slide, "位置来源：ExploreLoadingScreen.tscn 的全局 Control 矩形。")
    add_function_annotations(slide, title, scene, [], "加载进度和文案由探索房间加载流程更新；该界面无按钮。")


def slide_player_status(prs, title="玩家状态 HUD", scene="scenes/ui/player_status/PlayerStatusHUD.tscn"):
    slide = new_slide(prs, title, scene, "0 个交互点")
    layout = compute_layout("scenes/ui/player_status/PlayerStatusHUD.tscn")
    for name, label in [
        ("ScorePanel", "分数底板"),
        ("ScoreLabel", "分数"),
        ("YellowBar", "HP 黄条"),
        ("RedBar", "HP 红条"),
        ("BackBar", "狂热底条"),
        ("FillBar", "狂热填充"),
        ("Frame", "狂热边框"),
        ("Label", "狂热文字"),
    ]:
        draw_rect_from_layout(slide, layout, name, label)
    add_mapping_note(slide, "位置来源：PlayerStatusHUD.tscn，包含 Node2D 缩放后的条形子节点。")
    add_function_annotations(slide, title, scene, [], "显示 HP、分数、狂热条；该 HUD 不提供鼠标交互。")


def slide_hud_wrapper(prs):
    scene = "scenes/ui/hud.tscn"
    title = "HUD 画布包装"
    slide = new_slide(prs, title, scene, "0 个交互点")
    add_display_px(slide, 0, 0, 1920, 1080, "CanvasLayer 视口")
    layout = compute_layout("scenes/ui/player_status/PlayerStatusHUD.tscn")
    for name, label in [
        ("ScorePanel", "玩家 HUD 分数"),
        ("YellowBar", "玩家 HUD HP"),
        ("FillBar", "玩家 HUD 狂热"),
    ]:
        draw_rect_from_layout(slide, layout, name, label)
    add_mapping_note(slide, "hud.tscn 实例化 PlayerStatusHUD，位置沿用玩家 HUD 实际视口坐标。")
    add_function_annotations(slide, title, scene, [], "该场景是 HUD CanvasLayer 包装层，本身无按钮。")


def slide_boss_hud(prs):
    scene = "scenes/ui/BossHUD.tscn"
    title = "Boss HUD"
    slide = new_slide(prs, title, scene, "0 个交互点")
    layout = compute_layout(scene)
    for name, label in [
        ("FlashBar", "HP 闪烁条"),
        ("RedBar", "HP 红条"),
        ("NamePlate", "名称底板"),
        ("NameLabel", "Boss 名称"),
        ("Frame", "边框"),
    ]:
        draw_rect_from_layout(slide, layout, name, label)
    add_mapping_note(slide, "位置来源：BossHUD.tscn，包含 Control 缩放值。")
    add_function_annotations(slide, title, scene, [], "显示 Boss 名称与血条；该 HUD 不提供鼠标交互。")


def build():
    prs = Presentation()
    prs.slide_width = inch(SLIDE_W)
    prs.slide_height = inch(SLIDE_H)

    slide_main_menu(prs)
    slide_settings(prs)
    slide_boss_select(prs)
    slide_game_over(prs)
    slide_world_map(prs)
    slide_shop(prs)
    slide_hangar(prs)
    slide_equipment_row(prs)
    slide_event_result(prs)
    slide_evacuation(prs)
    slide_command_console(prs)
    slide_explore_map(prs)
    slide_compass(prs)
    slide_loading(prs)
    slide_player_status(prs)
    slide_hud_wrapper(prs)
    slide_boss_hud(prs)

    for slide in prs.slides:
        add_game_window_guide(slide)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)


if __name__ == "__main__":
    build()
